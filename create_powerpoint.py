#!/usr/bin/env python3
"""
Create a PowerPoint presentation for the SQL Server Upgrade Solution
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def format_slide_content(text_frame, font_size=16, space_after=4):
    """Helper function to format slide content consistently"""
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.space_after = Pt(space_after)
        paragraph.alignment = PP_ALIGN.LEFT
    
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.word_wrap = True

def create_sql_upgrade_presentation():
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SQL Server 2022 Upgrade Solution"
    subtitle.text = "Side-by-Side Instance Migration\nwith Enhanced Backup/Restore Options\n\nDatabase Administration Team"
    
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    
    format_slide_content(subtitle.text_frame, font_size=18, space_after=6)
    
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    content.text = """Complete SQL Server Instance Migration Solution

Key Benefits:
• Zero Downtime - Side-by-side upgrade approach
• Comprehensive Migration - Databases + Server Objects  
• Multiple Migration Methods - Direct, Backup/Restore, Detach/Attach
• Flexible Backup Options - New backups or existing backup chains
• Production Ready - Tested, validated, and documented
• Risk Mitigation - WhatIf mode and comprehensive logging"""
    
    format_slide_content(content.text_frame, font_size=16, space_after=4)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Architecture Overview"
    content.text = """Modular PowerShell Solution

6 Specialized Modules:
1. SQLUpgrade.Logging - Comprehensive logging
2. SQLUpgrade.Connection - Robust connection management
3. SQLUpgrade.Database - Database discovery and validation
4. SQLUpgrade.Encryption - TDE and encryption support
5. SQLUpgrade.Migration - Core migration functionality
6. SQLUpgrade.PostUpgrade - Post-migration tasks

Main Orchestrator: Start-SQLServerUpgrade.ps1
• Zero function definitions (calls modules only)
• Follows dbatools design patterns"""
    
    format_slide_content(content.text_frame, font_size=14, space_after=3)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Migration Methods"
    content.text = """Three Flexible Approaches

1. Direct Method (Default)
   • Uses Copy-DbaDatabase with automatic fallback
   • Best for most scenarios
   • Handles complex objects automatically

2. Backup/Restore Method
   • New Backups: Creates fresh full backup and restores
   • Existing Backups: Uses existing backup chains
   • Perfect for large databases

3. Detach/Attach Method
   • File-level database migration
   • Fastest for very large databases
   • Requires careful file path management"""
    
    format_slide_content(content.text_frame, font_size=14, space_after=3)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Enhanced Backup/Restore Features"
    content.text = """Comprehensive Backup Chain Support

New Backup Creation:
-MigrationMethod BackupRestore -BackupPath "C:\\Backups"

Existing Backup Chain:
-MigrationMethod BackupRestore -UseExistingBackups 
-FullBackupPath "C:\\Backups\\DB_Full.bak"
-DifferentialBackupPath "C:\\Backups\\DB_Diff.bak" 
-LogBackupPaths @("C:\\Backups\\Log1.trn", "Log2.trn")

Benefits:
• Leverage existing backup strategies
• Point-in-time recovery capability
• Minimal storage requirements
• Integration with backup policies"""
    
    format_slide_content(content.text_frame, font_size=13, space_after=2)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Server Object Migration"
    content.text = """Complete Instance Migration

11 Configurable Server Object Types:
• SQL Server Logins (excluding system)
• SQL Server Agent Jobs
• Linked Servers
• Server-level Triggers
• Custom Server Roles
• Credentials
• Proxy Accounts
• Alerts and Operators
• Backup Devices
• Server Configuration Settings

Flexible Selection:
• Individual switches: -IncludeLogins -IncludeJobs
• All objects: -IncludeAllServerObjects"""
    
    format_slide_content(content.text_frame, font_size=13, space_after=2)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "PowerShell Script Generation"
    content.text = """Offline Execution Capability

Generate PowerShell Scripts for Later Execution:
-OutputFile "C:\\Scripts\\Migration_Script.ps1"

Generated Scripts Include:
• dbatools-based migration commands
• Server object creation scripts
• Post-migration tasks (compatibility, statistics, indexes)
• Database integrity validation
• Comprehensive documentation and comments

Use Cases:
• Change management approval processes
• Scheduled maintenance windows
• Audit trail requirements
• Manual review and customization"""
    
    format_slide_content(content.text_frame, font_size=13, space_after=2)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Safety and Validation Features"
    content.text = """Enterprise-Grade Risk Management

WhatIf Mode:
• Preview all operations without changes
• Validate connections and compatibility
• Estimate migration scope and time
• Perfect for planning and approval

Idempotent Operations:
• Safe to run multiple times
• Skips existing objects automatically
• No data loss or corruption risk

Never Drops Objects:
• Only creates and migrates
• Preserves existing target data
• Additive approach only"""
    
    format_slide_content(content.text_frame, font_size=14, space_after=3)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Testing and Validation"
    content.text = """Production-Ready Quality Assurance

Comprehensive Test Suite:
• 61 Pester Tests - 100% pass rate
• Unit Tests - Individual module validation
• Integration Tests - End-to-end workflow testing
• Container Testing - Real SQL Server validation

Test Coverage:
• Module imports and function exports
• Parameter validation and error handling
• Connection management and security
• Migration methods and backup/restore
• WhatIf functionality and logging"""
    
    format_slide_content(content.text_frame, font_size=13, space_after=2)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Usage Examples"
    content.text = """Real-World Scenarios

Complete Instance Migration:
.\\Start-SQLServerUpgrade.ps1 
  -SourceInstance "SQL2019\\PROD" 
  -TargetInstance "SQL2022\\PROD" 
  -Databases "All" 
  -IncludeAllServerObjects

Existing Backups Migration:
.\\Start-SQLServerUpgrade.ps1 
  -MigrationMethod BackupRestore 
  -UseExistingBackups 
  -FullBackupPath "\\\\BackupServer\\DB_Full.bak"

Generate PowerShell Script:
.\\Start-SQLServerUpgrade.ps1 -OutputFile "Migration.ps1" -WhatIf"""
    
    format_slide_content(content.text_frame, font_size=12, space_after=2)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation Roadmap"
    content.text = """Deployment Strategy

Phase 1: Preparation (Week 1)
• Install PowerShell 7+ and dbatools module
• Review and customize configuration parameters
• Test connectivity to source and target instances
• Run WhatIf mode for migration planning

Phase 2: Pilot Testing (Week 2)
• Select non-critical databases for pilot migration
• Execute migrations in test environment
• Validate results and performance

Phase 3: Production Migration (Week 3-4)
• Schedule maintenance windows
• Execute production migrations
• Monitor and validate results"""
    
    format_slide_content(content.text_frame, font_size=13, space_after=2)
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Questions and Next Steps"
    content.text = """Ready for Implementation

Solution Benefits Recap:
• Complete Instance Migration - Databases + Server Objects
• Multiple Migration Methods - Flexible approach selection
• Enhanced Backup/Restore - Existing backup chain support
• PowerShell Script Generation - Offline execution capability
• Enterprise Safety Features - WhatIf, logging, validation
• Production Tested - 100% test pass rate

Repository Location:
• GitHub: karim-attaleb/sql-server-instance-upgrade
• Branch: sql-server-upgrade-solution

Questions?
• Technical implementation details
• Specific migration scenarios
• Timeline and resource planning"""
    
    format_slide_content(content.text_frame, font_size=13, space_after=2)
    
    prs.save('/home/ubuntu/sql-server-instance-upgrade/SQL-Server-Upgrade-Solution-Presentation.pptx')
    print("PowerPoint presentation created successfully!")

if __name__ == "__main__":
    create_sql_upgrade_presentation()
